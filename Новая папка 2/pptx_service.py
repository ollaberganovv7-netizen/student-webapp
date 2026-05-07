import io
import os
import copy
from typing import List
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt, Emu


def analyze_template(template_path: str) -> dict:
    """
    Analyze a PPTX template and return detailed info about each slide's structure.
    Reports per-block word estimates so AI can generate properly sized content.
    
    Returns:
        {"total_slides": int, "slides": [{"index": int, "text_areas": int, 
         "has_title": bool, "has_body": bool, "body_blocks": int,
         "blocks": [{"words": int, "type": "body"|"textbox"}, ...],
         "estimated_words": int}, ...]}
    """
    if template_path and not os.path.isabs(template_path):
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        full = os.path.join(base, template_path)
        if os.path.exists(full):
            template_path = full

    if not template_path or not os.path.exists(template_path):
        return {"total_slides": 0, "slides": []}

    prs = Presentation(template_path)
    result = {"total_slides": len(prs.slides), "slides": []}
    
    for i, slide in enumerate(prs.slides):
        info = {"index": i, "text_areas": 0, "has_title": False, 
                "has_body": False, "body_blocks": 0, "blocks": [],
                "estimated_words": 0}
        
        # Sort shapes geographically (left-to-right, then top-to-bottom)
        # We round X to nearest ~0.5 inch (457200 EMU) to group columns together
        sorted_shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame],
            key=lambda s: (round(getattr(s, 'left', 0) / 457200), getattr(s, 'top', 0))
        )
        
        for shape in sorted_shapes:
            info["text_areas"] += 1
            if shape.is_placeholder:
                pf = shape.placeholder_format
                if pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    info["has_title"] = True
                elif pf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    info["has_body"] = True
                    info["body_blocks"] += 1
                    # Estimate capacity from placeholder size
                    w_inches = shape.width / 914400  # EMU to inches
                    h_inches = shape.height / 914400
                    area = w_inches * h_inches
                    block_words = max(15, int(area * 10))
                    info["blocks"].append({"words": block_words, "type": "body"})
                    info["estimated_words"] += block_words
            else:
                # Regular text box
                w_inches = shape.width / 914400
                h_inches = shape.height / 914400
                area = w_inches * h_inches
                block_words = max(10, int(area * 8))
                info["blocks"].append({"words": block_words, "type": "textbox"})
                info["estimated_words"] += block_words
        
        result["slides"].append(info)
    
    return result

# ── Slide shape analysis ─────────────────────────────────────────────────────

TITLE_PH = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_PH  = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
SUB_PH   = {PP_PLACEHOLDER.SUBTITLE}

def _classify_shapes(slide):
    """Return dict of shape lists: title, body, subtitle, textboxes."""
    title, body, subtitle, textboxes = [], [], [], []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            pf = shape.placeholder_format
            if pf.type in TITLE_PH:
                title.append(shape)
            elif pf.type in BODY_PH:
                body.append(shape)
            elif pf.type in SUB_PH:
                subtitle.append(shape)
            else:
                textboxes.append(shape)
        else:
            textboxes.append(shape)
    return {"title": title, "body": body, "subtitle": subtitle, "other": textboxes}


# ── Smart Template Engine (auto-detect title/body in ANY template) ────────────

def _get_font_size_from_shape(shape):
    """Get the largest font size used in a shape (in Pt)."""
    max_size = 0
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size.pt
                    if size_pt > max_size:
                        max_size = size_pt
    except Exception:
        pass
    return max_size


def _smart_classify_textboxes(textboxes, slide):
    """
    Intelligently classify regular text boxes as title or body
    based on position, font size, and dimensions.
    
    Rules:
    1. The shape closest to the TOP with the LARGEST font = Title
    2. The shape with the LARGEST area (excluding title) = Body
    3. Everything else = decorative (skip)
    """
    if not textboxes:
        return None, None
    
    # Filter out tiny decorative elements (less than 1.5 inches wide)
    candidates = []
    for s in textboxes:
        if s.width >= Inches(1.5) and s.height >= Inches(0.3):
            candidates.append(s)
    
    if not candidates:
        return None, None
    
    # Score each candidate for "title-ness":
    # Higher score = more likely to be the title
    # Factors: top position (higher = more title-like), big font, short text
    best_title = None
    best_title_score = -1
    
    for s in candidates:
        score = 0
        font_size = _get_font_size_from_shape(s)
        text_len = len(s.text_frame.text.strip())
        
        # Big font = likely title (font size contributes heavily)
        score += font_size * 3
        
        # Top of slide = likely title (invert top position so higher = better)
        # Slide height is typically ~6858000 EMU (7.5 inches)
        top_normalized = max(0, 6858000 - s.top) / 6858000  # 1.0 at top, 0.0 at bottom
        score += top_normalized * 100
        
        # Short text = likely title (penalize long text)
        if text_len < 80:
            score += 50
        elif text_len > 200:
            score -= 50
        
        if score > best_title_score:
            best_title_score = score
            best_title = s
    
    # Body = the LARGEST remaining text box (by area)
    best_body = None
    best_body_area = 0
    
    for s in candidates:
        if s is best_title:
            continue
        area = s.width * s.height
        if area > best_body_area:
            best_body_area = area
            best_body = s
    
    # If we only found one candidate, decide: is it a title or body?
    if best_body is None and best_title is not None:
        font_size = _get_font_size_from_shape(best_title)
        if font_size < 20:  # Small font = probably body, not title
            best_body = best_title
            best_title = None
    
    return best_title, best_body


# ── Format-preserving text fill ───────────────────────────────────────────────

def _capture_font(run):
    """Capture font properties from a Run."""
    f = run.font
    color_val = None
    color_type = None
    try:
        if getattr(f, 'color', None) and getattr(f.color, 'type', None):
            # python-pptx raises an error if you access .rgb on a scheme color
            if str(f.color.type).endswith('RGB') or f.color.type == 1: # MSO_COLOR_TYPE.RGB
                color_val = f.color.rgb
                color_type = "rgb"
            elif str(f.color.type).endswith('SCHEME') or f.color.type == 2: # MSO_COLOR_TYPE.SCHEME
                color_val = getattr(f.color, 'theme_color', None)
                color_type = "theme"
    except Exception:
        pass

    return {
        "size": getattr(f, 'size', None),
        "bold": getattr(f, 'bold', None),
        "italic": getattr(f, 'italic', None),
        "name": getattr(f, 'name', None),
        "color_val": color_val,
        "color_type": color_type,
    }


def _apply_font(font, props):
    """Apply captured font properties."""
    try:
        if props.get("size"):    font.size = props["size"]
        if props.get("bold") is not None: font.bold = props["bold"]
        if props.get("italic") is not None: font.italic = props["italic"]
        if props.get("name"):    font.name = props["name"]
        
        c_val = props.get("color_val")
        c_type = props.get("color_type")
        if c_val is not None:
            if c_type == "rgb":
                font.color.rgb = c_val
            elif c_type == "theme":
                font.color.theme_color = c_val
    except Exception:
        pass


def _fill_textframe(tf, content, is_title=False, compact=False):
    """Replace text in a TextFrame, preserving first-run formatting.
    compact=True: used on picture-layout slides, reduces font/word limits."""
    # capture formatting from first existing run
    font_props = {}
    alignment = None
    if tf.paragraphs:
        p0 = tf.paragraphs[0]
        alignment = p0.alignment
        if p0.runs:
            font_props = _capture_font(p0.runs[0])

    tf.clear()
    tf.word_wrap = True

    items = content if isinstance(content, list) else [content]
    
    # Limit body text to prevent overflow
    if not is_title:
        max_points = 2 if compact else 5  # Allow more points (AI controls count now)
        if len(items) > max_points:
            items = items[:max_points]

    # Set default and max font sizes
    if compact:
        default_size = Pt(24) if is_title else Pt(14)
        max_size = Pt(32) if is_title else Pt(16)
        max_words = 150
    else:
        default_size = Pt(28) if is_title else Pt(16)
        max_size = Pt(40) if is_title else Pt(20)
        max_words = 250

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text = str(item)
        # Trim very long points
        if not is_title:
            words = text.split()
            if len(words) > max_words:
                text = ' '.join(words[:max_words]) + '...'
        para.text = text
        para.level = 0
        if alignment is not None:
            para.alignment = alignment
        for run in para.runs:
            _apply_font(run.font, font_props)
            # Override font size to ensure readability, but clamp to max_size
            size = font_props.get("size") or default_size
            if size > max_size:
                size = max_size
            run.font.size = size

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


# ── Deep slide duplication (XML level) ────────────────────────────────────────

def _duplicate_slide(prs, source_slide):
    """
    Duplicate a slide preserving visual design elements.
    
    Uses add_slide(layout) for writable placeholders, then copies
    non-placeholder shapes (decorative elements, images, backgrounds)
    from the source slide.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Copy background from source if it exists
    src_bg = source_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
    )
    if src_bg is not None:
        new_bg = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if new_bg is not None:
            new_slide._element.remove(new_bg)
        # Insert bg before spTree
        spTree = new_slide._element.find(
            './/{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/'
            '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spTree'
        )
        cSld = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        if cSld is not None:
            cSld.insert(0, copy.deepcopy(src_bg))

    # Copy non-placeholder shapes (decorative elements, images, etc.)
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    src_spTree = source_slide._element.find(f'{{{ns_p}}}cSld/{{{ns_a}}}spTree')
    new_spTree = new_slide._element.find(f'{{{ns_p}}}cSld/{{{ns_a}}}spTree')
    
    if src_spTree is not None and new_spTree is not None:
        for child in src_spTree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            # Copy pics, connectors, group shapes (not sp placeholders)
            if tag in ('pic', 'cxnSp', 'grpSp', 'graphicFrame'):
                new_spTree.append(copy.deepcopy(child))
            elif tag == 'sp':
                # Only copy non-placeholder shapes (decorative text boxes, etc.)
                nvSpPr = child.find(f'.//{{{ns_p}}}nvSpPr/{{{ns_p}}}nvPr')
                if nvSpPr is None:
                    nvSpPr = child.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
                has_ph = False
                if nvSpPr is not None:
                    ph_elem = nvSpPr.find('{http://schemas.openxmlformats.org/presentationml/2006/main}ph')
                    has_ph = ph_elem is not None
                if not has_ph:
                    new_spTree.append(copy.deepcopy(child))

    # Copy image relationships from source
    for rel in source_slide.part.rels.values():
        if rel.is_external:
            continue
        try:
            if 'image' in rel.reltype.lower():
                new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            pass

    return new_slide


# ── Image helper ──────────────────────────────────────────────────────────────

def _find_picture_placeholder(slide):
    """Find picture placeholder in slide (the frame/рамка for images).
    Only detects actual PICTURE placeholders (type=18), not OBJECT/BODY."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            pf = shape.placeholder_format
            # PP_PLACEHOLDER.PICTURE = 18 — this is the real picture frame
            if pf.type == PP_PLACEHOLDER.PICTURE:
                return shape
            # Also BITMAP, MEDIA_CLIP
            if pf.type in (PP_PLACEHOLDER.BITMAP, PP_PLACEHOLDER.MEDIA_CLIP):
                return shape
    return None


def _add_image(slide, img_bytes, prs, position="right"):
    """Insert image into picture placeholder if available, else overlay at position."""
    try:
        # ── PHASE 1: Try to find a picture placeholder (рамка) ──
        pic_ph = _find_picture_placeholder(slide)
        if pic_ph is not None:
            try:
                stream = io.BytesIO(img_bytes)
                pic_ph.insert_picture(stream)
                return  # Successfully inserted into placeholder frame
            except Exception:
                pass  # Placeholder doesn't support insert_picture, fall through

        # ── PHASE 2: Fallback — overlay image at calculated position ──
        stream = io.BytesIO(img_bytes)
        w, h = prs.slide_width, prs.slide_height
        if position == "center":
            iw, ih = int(w * 0.35), int(h * 0.35)
            slide.shapes.add_picture(stream, int((w - iw) / 2), int(h * 0.56), iw, ih)
        else:
            iw, ih = int(w * 0.35), int(h * 0.55)
            slide.shapes.add_picture(stream, int(w - iw - Inches(0.3)),
                                     int((h - ih) / 2), iw, ih)
    except Exception as e:
        print(f"Image insert error: {e}")



# ── Fill a single slide with AI data ──────────────────────────────────────────

def _fill_slide(slide, slide_info):
    """Fill title + body on any slide."""
    title_text = slide_info.get("title", "")
    points = slide_info.get("points", [])
    shapes = _classify_shapes(slide)
    
    # Check if this slide has a picture placeholder → compact mode
    has_picture = _find_picture_placeholder(slide) is not None
    if has_picture and len(points) > 2:
        # Trim to 2 shorter points for picture-layout slides
        points = points[:2]
    
    has_placeholders = bool(shapes["title"] or shapes["body"] or shapes["subtitle"])

    # ── Fill title placeholder(s) ──
    title_filled = False
    for s in shapes["title"]:
        _fill_textframe(s.text_frame, title_text, is_title=True)
        title_filled = True

    # ── Fill body placeholder(s) — ONLY real placeholders ──
    body_filled = False
    if shapes["body"]:
        # Distribute points across all available body placeholders
        num_bodies = len(shapes["body"])
        points_per_body = max(1, len(points) // num_bodies)
        
        for i, s in enumerate(shapes["body"]):
            start_idx = i * points_per_body
            # If it's the last body, give it all the remaining points
            end_idx = (i + 1) * points_per_body if i < num_bodies - 1 else len(points)
            
            chunk = points[start_idx:end_idx]
            if chunk:
                _fill_textframe(s.text_frame, chunk, compact=has_picture)
                body_filled = True
            else:
                # If we ran out of points, just clear the placeholder so it doesn't show "Текст слайда"
                s.text_frame.clear()

    # ── Fallback: subtitle placeholder ──
    if not body_filled and points:
        for s in shapes["subtitle"]:
            _fill_textframe(s.text_frame, points)
            body_filled = True
            break

    # ── SMART FALLBACK: Auto-detect from regular text boxes ──
    if not has_placeholders and shapes["other"]:
        smart_title, smart_body = _smart_classify_textboxes(shapes["other"], slide)
        
        if smart_title and not title_filled:
            _fill_textframe(smart_title.text_frame, title_text, is_title=True)
            title_filled = True
        
        if smart_body and not body_filled and points:
            _fill_textframe(smart_body.text_frame, points)
            body_filled = True
    
    # ── Last resort: single largest "other" textbox ──
    if not body_filled and points and shapes["other"]:
        candidates = []
        for s in shapes["other"]:
            if s.width < Inches(1.5) or s.height < Inches(0.5):
                continue
            candidates.append(s)
        
        if candidates:
            biggest = max(candidates, key=lambda s: s.width * s.height)
            _fill_textframe(biggest.text_frame, points)
            body_filled = True


# ── Layout fallback (for templates with only 1 slide) ─────────────────────────

def _best_content_layout(prs):
    preferred = ["Заголовок и объект", "Title and Content", "Sarlavha va ob'ekt"]
    for name in preferred:
        for lay in prs.slide_layouts:
            if name.lower() in lay.name.lower():
                return lay
    for lay in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in lay.placeholders}
        if types & TITLE_PH and types & BODY_PH:
            return lay
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


# ── Remove extra slides ──────────────────────────────────────────────────────

def _remove_slide(prs, slide_index):
    """Remove slide at given index."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if slide_index < len(items):
        el = items[slide_index]
        try:
            prs.part.drop_rel(el.rId)
        except Exception:
            pass
        sldIdLst.remove(el)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def generate_pptx(
    slides_data: List[dict],
    template_path: str,
    topic: str,
    author: str = "Foydalanuvchi",
    slide_images: dict = None,
    maket_mode: bool = False,
) -> bytes:
    """
    Smart PPTX generator.

    • If the template has multiple pre-designed slides → fills them in-place,
      duplicating the last content slide when more are needed.
    • If the template has only a title slide → falls back to adding slides
      from the best matching layout (old behaviour).
    """

    # resolve path
    if template_path and not os.path.isabs(template_path):
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        full = os.path.join(base, template_path)
        if os.path.exists(full):
            template_path = full

    prs = (Presentation(template_path)
           if template_path and os.path.exists(template_path)
           else Presentation())
    slide_images = slide_images or {}

    tpl_count = len(prs.slides)          # total slides in template
    tpl_content = max(0, tpl_count - 1)  # content slides (minus title)
    
    # In maket_mode, we only fill as many slides as exist in the template
    if maket_mode:
        need_content = min(len(slides_data), tpl_content)
    else:
        need_content = len(slides_data)      # how many AI content slides

    # ── 1. Title slide ────────────────────────────────────────────────────
    if tpl_count > 0:
        ts = prs.slides[0]

        # Smart Engine fallback for title slide
        shapes = _classify_shapes(ts)
        title_filled = False
        subtitle_filled = False
        
        for s in shapes["title"]:
            _fill_textframe(s.text_frame, topic, is_title=True)
            title_filled = True
        for s in shapes["subtitle"]:
            _fill_textframe(s.text_frame, f"Bajardi: {author}")
            subtitle_filled = True
        
        # Smart fallback for title slide (Canva templates without placeholders)
        if not title_filled and shapes["other"]:
            smart_title, smart_body = _smart_classify_textboxes(shapes["other"], ts)
            if smart_title:
                _fill_textframe(smart_title.text_frame, topic, is_title=True)
            if smart_body:
                _fill_textframe(smart_body.text_frame, f"Bajardi: {author}")
        
        if 0 in slide_images:
            _add_image(ts, slide_images[0], prs, "center")

    # ── 2. Content slides ─────────────────────────────────────────────────
    if tpl_content == 0:
        # --- Template has only title → create from layout (old way) ---
        layout = _best_content_layout(prs)
        for i, info in enumerate(slides_data):
            ns = prs.slides.add_slide(layout)
            _fill_slide(ns, info)
            if (i + 1) in slide_images:
                _add_image(ns, slide_images[i + 1], prs)
    else:
        # --- Template has pre-designed content slides → use them ---

        # Classify template content slides: which have picture placeholders?
        tpl_slides = list(prs.slides)
        text_only_indices = []
        picture_indices = []

        for si in range(1, tpl_count):
            slide = tpl_slides[si]
            if _find_picture_placeholder(slide) is not None:
                picture_indices.append(si)
            else:
                text_only_indices.append(si)

        # Step A: Duplicate all needed content slides FIRST (sources still alive)
        if not maket_mode:
            text_cycle = 0
            for i in range(need_content):
                slide_idx = i + 1  # slide index in final presentation
                has_image = slide_idx in slide_images
                
                if has_image and picture_indices:
                    # Use the first picture layout (or could cycle if multiple)
                    _duplicate_slide(prs, tpl_slides[picture_indices[0]])
                else:
                    if text_only_indices:
                        # Alternate between available text layouts for variety
                        idx_to_use = text_only_indices[text_cycle % len(text_only_indices)]
                        _duplicate_slide(prs, tpl_slides[idx_to_use])
                        text_cycle += 1
                    elif picture_indices:
                        _duplicate_slide(prs, tpl_slides[picture_indices[0]])
                    else:
                        _duplicate_slide(prs, tpl_slides[tpl_count - 1])

        # Step B: Remove the ORIGINAL template content slides (they come before duplicates)
        # Original content slides are at indices 1..tpl_count-1
        # Remove them from last to first to avoid index shifting
        if not maket_mode:
            for _ in range(tpl_content):
                _remove_slide(prs, 1)  # always remove index 1 (first content slide)

        # Step C: Fill all content slides
        fresh = list(prs.slides)
        for i in range(need_content):
            idx = i + 1  # skip title
            if idx < len(fresh):
                _fill_slide(fresh[idx], slides_data[i])
                if (i + 1) in slide_images:
                    _add_image(fresh[idx], slide_images[i + 1], prs)

    # ── 3. Save ───────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()